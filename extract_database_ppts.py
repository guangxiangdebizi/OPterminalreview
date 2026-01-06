"""
数据库期末复习PPT内容提取与整理脚本
按章节提取PPT文本并生成对应的复习资料
"""

import os
from pathlib import Path
from pptx import Presentation
import re
from collections import defaultdict

def extract_text_from_ppt(ppt_path):
    """从PPT中提取文本内容"""
    try:
        prs = Presentation(ppt_path)
        slides_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_content.append(f"\n### 幻灯片 {slide_num}\n")
            
            # 提取所有文本框中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                
                # 如果是表格，提取表格内容
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(" | ".join(row_data))
                    if table_data:
                        slide_content.append("\n**表格内容：**\n")
                        slide_content.append("\n".join(table_data))
            
            if len(slide_content) > 1:  # 如果有内容（除了标题）
                slides_text.append("\n".join(slide_content))
        
        return "\n\n".join(slides_text)
    except Exception as e:
        print(f"处理 {ppt_path} 时出错: {str(e)}")
        return ""

def get_chapter_number(filename):
    """从文件名中提取章节号"""
    # 匹配各种章节号格式
    patterns = [
        r'^(\d+)\..*',  # 匹配 "2.1", "3.1" 等
        r'^(\d+)\s',  # 匹配 "5 触发器" 等
        r'^(\d+)第',  # 匹配 "6第六章" 等
        r'第([一二三四五六七八九十]+)章',  # 匹配 "第二章", "第三章" 等
        r'第(\d+)章',  # 匹配 "第2章", "第3章" 等
    ]
    
    for pattern in patterns:
        match = re.search(pattern, filename)
        if match:
            chapter = match.group(1)
            # 转换中文数字为阿拉伯数字
            chinese_to_arabic = {
                '一': '1', '二': '2', '三': '3', '四': '4', '五': '5',
                '六': '6', '七': '7', '八': '8', '九': '9', '十': '10'
            }
            if chapter in chinese_to_arabic:
                return chinese_to_arabic[chapter]
            return chapter.split('.')[0]  # 如果是 "2.1"，只取 "2"
    
    return "0"  # 默认章节

def main():
    # 数据库PPT文件夹路径
    ppt_folder = Path("数据库期末复习")
    
    if not ppt_folder.exists():
        print(f"错误：找不到文件夹 {ppt_folder}")
        return
    
    # 获取所有pptx文件
    ppt_files = list(ppt_folder.glob("*.pptx"))
    
    if not ppt_files:
        print("未找到任何PPT文件")
        return
    
    print(f"找到 {len(ppt_files)} 个PPT文件")
    
    # 按章节组织内容
    chapters = defaultdict(list)
    
    # 处理每个PPT文件
    for ppt_file in ppt_files:
        filename = ppt_file.name
        print(f"\n正在处理: {filename}")
        
        # 获取章节号
        chapter_num = get_chapter_number(filename)
        
        # 提取文本
        text_content = extract_text_from_ppt(ppt_file)
        
        if text_content:
            chapters[chapter_num].append({
                'filename': filename,
                'content': text_content
            })
            print(f"  - 归类到第 {chapter_num} 章")
    
    # 输出文件夹
    output_folder = Path("数据库期末复习资料")
    output_folder.mkdir(exist_ok=True)
    
    # 章节名称映射
    chapter_names = {
        '2': '关系数据库',
        '3': 'SQL语言',
        '4': '存储过程与编程',
        '5': '触发器',
        '6': '索引及查询优化',
        '7': '关系数据库设计理论',
        '8': '数据库设计',
        '9': '数据库安全',
        '10': '数据库保护'
    }
    
    # 生成每个章节的复习资料
    for chapter_num in sorted(chapters.keys(), key=lambda x: int(x) if x.isdigit() else 0):
        chapter_name = chapter_names.get(chapter_num, f'第{chapter_num}章')
        output_file = output_folder / f"数据库第{chapter_num}章_{chapter_name}.md"
        
        print(f"\n生成第 {chapter_num} 章复习资料...")
        
        with open(output_file, 'w', encoding='utf-8') as f:
            # 写入标题
            f.write(f"# 数据库第{chapter_num}章 - {chapter_name}\n\n")
            f.write(f"> 整理自数据库期末复习PPT\n\n")
            f.write("---\n\n")
            
            # 写入每个PPT的内容
            for ppt_info in chapters[chapter_num]:
                f.write(f"## {ppt_info['filename']}\n\n")
                f.write(ppt_info['content'])
                f.write("\n\n---\n\n")
        
        print(f"  ✓ 已生成: {output_file}")
    
    # 生成总目录
    index_file = output_folder / "README.md"
    with open(index_file, 'w', encoding='utf-8') as f:
        f.write("# 数据库期末复习资料\n\n")
        f.write("> 从PPT中提取并整理的章节复习资料\n\n")
        f.write("## 目录\n\n")
        
        for chapter_num in sorted(chapters.keys(), key=lambda x: int(x) if x.isdigit() else 0):
            chapter_name = chapter_names.get(chapter_num, f'第{chapter_num}章')
            filename = f"数据库第{chapter_num}章_{chapter_name}.md"
            f.write(f"- [第{chapter_num}章 - {chapter_name}](./{filename})\n")
            
            # 列出该章节包含的PPT文件
            for ppt_info in chapters[chapter_num]:
                f.write(f"  - {ppt_info['filename']}\n")
    
    print(f"\n✓ 所有复习资料已生成到文件夹: {output_folder}")
    print(f"✓ 共生成 {len(chapters)} 个章节的复习资料")

if __name__ == "__main__":
    main()

